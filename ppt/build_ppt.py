#!/usr/bin/env python3
"""
Build the hackathon presentation PowerPoint.
Run from project root:
  uv run --with python-pptx --with matplotlib --with pandas --with scipy --with numpy python3 ppt/build_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
from pathlib import Path

OUT = Path("ppt/Siglec6_Hackathon.pptx")

# ── palette ────────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1A, 0x1A, 0x2E)   # near-black navy
MID_BG     = RGBColor(0x16, 0x21, 0x3E)   # dark blue
ACCENT1    = RGBColor(0x0F, 0x3A, 0x67)   # deep blue
BOLTZ_CLR  = RGBColor(0x2E, 0x86, 0xAB)   # teal-blue   (Boltz-2)
AF3A_CLR   = RGBColor(0xA2, 0x3B, 0x72)   # magenta     (AF3-MSA)
AF3B_CLR   = RGBColor(0xF1, 0x8F, 0x01)   # amber       (AF3-ColabFold)
SCORE_CLR  = RGBColor(0x2D, 0x93, 0xAD)   # cyan
VINA_CLR   = RGBColor(0x44, 0xBF, 0x6E)   # green
KD_CLR     = RGBColor(0xE8, 0x3A, 0x3A)   # red
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ARROW_CLR  = RGBColor(0x88, 0x88, 0xAA)
GOLD_CLR   = RGBColor(0xFF, 0xCC, 0x00)

SW = Inches(13.33)   # slide width  (widescreen 16:9)
SH = Inches(7.5)     # slide height


# ── helpers ────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h,
             fill: RGBColor | None = None,
             line: RGBColor | None = None,
             line_width: Pt = Pt(1.5),
             radius: bool = True) -> object:
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE = 1; rounded = 5
        x, y, w, h
    )
    sf = shape.fill
    if fill:
        sf.solid(); sf.fore_color.rgb = fill
    else:
        sf.background()
    ln = shape.line
    if line:
        ln.color.rgb = line; ln.width = line_width
    else:
        ln.fill.background()
    return shape


def add_rounded_rect(slide, x, y, w, h,
                     fill: RGBColor, line: RGBColor | None = None,
                     line_width=Pt(1.5)):
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        5,   # msoShapeRoundedRectangle
        x, y, w, h
    )
    shape.adjustments[0] = 0.08   # corner radius
    sf = shape.fill; sf.solid(); sf.fore_color.rgb = fill
    ln = shape.line
    if line:
        ln.color.rgb = line; ln.width = line_width
    else:
        ln.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=Pt(11), bold=False, color=WHITE,
             align=PP_ALIGN.CENTER, wrap=True) -> object:
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def add_label_in_shape(slide, shape, text, font_size=Pt(10), bold=False, color=WHITE):
    """Add centered text over an existing shape by placing a transparent textbox on top."""
    txb = slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
    tf  = txb.text_frame
    tf.word_wrap = True
    from pptx.util import Pt as _Pt
    from pptx.enum.text import PP_ALIGN
    # vertical centering via margin
    tf.margin_top    = shape.height // 4
    tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.size  = font_size
        run.font.bold  = bold
        run.font.color.rgb = color


def add_arrow(slide, x1, y1, x2, y2, color=ARROW_CLR, width=Pt(2)):
    """Simple line-with-arrowhead connector."""
    from pptx.util import Emu
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)   # 1=straight
    connector.line.color.rgb = color
    connector.line.width     = width
    # arrowhead at end
    from pptx.oxml.ns import qn
    from lxml import etree
    ln_elem = connector.line._ln
    tail_end = etree.SubElement(ln_elem, qn("a:tailEnd"))
    head_end = ln_elem.find(qn("a:headEnd"))
    if head_end is None:
        head_end = etree.SubElement(ln_elem, qn("a:headEnd"))
    head_end.set("type", "none")
    tail_end_elem = ln_elem.find(qn("a:tailEnd"))
    tail_end_elem.set("type", "triangle")
    tail_end_elem.set("w", "med")
    tail_end_elem.set("len", "med")


def slide_header_footer(slide, title_text, subtitle_text=""):
    """Add standard dark slide top/bottom bars and title."""
    add_rect(slide, 0, 0, SW, Inches(0.08), fill=BOLTZ_CLR)
    add_rect(slide, 0, SH - Inches(0.08), SW, Inches(0.08), fill=AF3A_CLR)
    add_text(slide, title_text,
             Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.60),
             font_size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle_text:
        add_text(slide, subtitle_text,
                 Inches(0.4), Inches(0.68), Inches(12.5), Inches(0.35),
                 font_size=Pt(10), color=LIGHT_GRAY, align=PP_ALIGN.LEFT)


def embed_figure(slide, img_path, x, y, w, h):
    """Embed a PNG figure into a slide."""
    slide.shapes.add_picture(str(img_path), x, y, w, h)


def add_caption(slide, text, y=None):
    """Add small caption text at the bottom of the slide."""
    if y is None:
        y = SH - Inches(0.42)
    add_text(slide, text,
             Inches(0.4), y, Inches(12.5), Inches(0.35),
             font_size=Pt(9), color=LIGHT_GRAY, align=PP_ALIGN.LEFT)


# ── Slide 1: Title ─────────────────────────────────────────────────────────

def slide_title(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, DARK_BG)

    # Top accent bar
    bar = add_rect(slide, 0, 0, SW, Inches(0.08), fill=BOLTZ_CLR)

    # Institution / event line
    add_text(slide, "Scripps Research  ·  CBB Hackathon 2026  ·  May 29 – June 1",
             Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.4),
             font_size=Pt(11), color=LIGHT_GRAY)

    # Main title
    add_text(slide,
             "In Silico Screening of Siglec-6 Ligands\nUsing Multi-Method Cofolding & Docking",
             Inches(1), Inches(1.3), Inches(11.3), Inches(2.2),
             font_size=Pt(36), bold=True, color=WHITE)

    # Subtitle / methods line
    add_text(slide,
             "Boltz-2  ·  AlphaFold3  ·  AutoDock Vina  ·  Experimental Kd Validation",
             Inches(1), Inches(3.3), Inches(11.3), Inches(0.6),
             font_size=Pt(16), color=LIGHT_GRAY)

    # Authors — side by side
    add_text(slide, "Suraj Ugrani",
             Inches(1), Inches(5.8), Inches(5.5), Inches(0.5),
             font_size=Pt(13), color=LIGHT_GRAY, align=PP_ALIGN.LEFT)
    add_text(slide, "Julien Heberling",
             Inches(6.8), Inches(5.8), Inches(5.5), Inches(0.5),
             font_size=Pt(13), color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

    # Bottom bar
    add_rect(slide, 0, SH - Inches(0.08), SW, Inches(0.08), fill=AF3A_CLR)
    return slide


# ── Slide 2: Pipeline Schematic ────────────────────────────────────────────

def slide_pipeline(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, SW, Inches(0.08), fill=BOLTZ_CLR)
    add_rect(slide, 0, SH - Inches(0.08), SW, Inches(0.08), fill=AF3A_CLR)

    # Title
    add_text(slide, "Computational Pipeline Overview",
             Inches(0.3), Inches(0.12), Inches(12.7), Inches(0.55),
             font_size=Pt(20), bold=True, color=WHITE)

    # ── Row 1: Inputs ──────────────────────────────────────────────────────
    row1_y = Inches(0.85)
    bh     = Inches(0.75)   # box height

    # Protein box
    p_box = add_rounded_rect(slide, Inches(0.4), row1_y, Inches(3.2), bh, fill=ACCENT1)
    add_label_in_shape(slide, p_box,
                       "Siglec-6 Protein  (FASTA)\n333 residues  ·  Vtype + C2type1 + C2type2",
                       font_size=Pt(9), bold=False)

    # Ligand box
    l_box = add_rounded_rect(slide, Inches(4.0), row1_y, Inches(2.6), bh, fill=ACCENT1)
    add_label_in_shape(slide, l_box,
                       "14-Compound Library\n13 analogs + sialic acid control",
                       font_size=Pt(9))

    # MSA box
    m_box = add_rounded_rect(slide, Inches(7.0), row1_y, Inches(2.6), bh, fill=ACCENT1)
    add_label_in_shape(slide, m_box,
                       "Pre-computed MSA\n(ColabFold  ·  .a3m format)",
                       font_size=Pt(9))

    # Arrow inputs → cofolding
    mid1_x = Inches(0.4) + Inches(3.2) // 2
    mid2_x = Inches(4.0) + Inches(2.6) // 2
    mid3_x = Inches(7.0) + Inches(2.6) // 2
    row2_y = Inches(2.0)

    for mx in [mid1_x, mid2_x, mid3_x]:
        add_arrow(slide, mx, row1_y + bh, mx, row2_y, color=ARROW_CLR)

    # ── Row 2: Cofolding label ─────────────────────────────────────────────
    add_text(slide, "COFOLDING  (protein + ligand structure prediction)",
             Inches(0.3), Inches(1.75), Inches(10), Inches(0.28),
             font_size=Pt(8.5), bold=True, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

    # ── Row 2: 3 cofolding boxes ───────────────────────────────────────────
    row2_y = Inches(2.0)
    cbh    = Inches(1.1)
    cbw    = Inches(3.0)
    gaps   = [Inches(0.4), Inches(4.0), Inches(7.6)]

    methods = [
        (BOLTZ_CLR,  "Boltz-2",
         "Open-source  ·  AWS EC2 g5.xlarge\nDocker  ·  GPU  ·  ColabFold MSA input"),
        (AF3A_CLR,   "AlphaFold3  (AF3-MSA)",
         "AF3 server  ·  built-in MSA search\nAutomated MSA pipeline"),
        (AF3B_CLR,   "AlphaFold3  (ColabFold-MSA)",
         "AF3 server  ·  pre-computed MSA\nSame MSA as Boltz-2 input"),
    ]

    method_centers = []
    for (clr, title, sub), gx in zip(methods, gaps):
        box = add_rounded_rect(slide, gx, row2_y, cbw, cbh, fill=clr)
        add_label_in_shape(slide, box, f"{title}\n{sub}", font_size=Pt(8.5), bold=False)
        method_centers.append(gx + cbw // 2)

    # ── Row 3: "42 structures" summary ────────────────────────────────────
    row3_y = Inches(3.3)
    sumbox = add_rounded_rect(slide, Inches(3.5), row3_y, Inches(4.5), Inches(0.65),
                              fill=RGBColor(0x2A, 0x2A, 0x4A),
                              line=LIGHT_GRAY, line_width=Pt(1))
    add_label_in_shape(slide, sumbox,
                       "42 Predicted Structures   (14 ligands  ×  3 methods)",
                       font_size=Pt(10), bold=True)

    # Arrows cofolding → summary
    sum_top_x = Inches(3.5) + Inches(4.5) // 2
    sum_top_y = row3_y
    for cx in method_centers:
        add_arrow(slide, cx, row2_y + cbh, sum_top_x, sum_top_y, color=ARROW_CLR)

    # ── Row 4: Analysis blocks ────────────────────────────────────────────
    add_text(slide, "ANALYSIS",
             Inches(0.3), Inches(4.1), Inches(3), Inches(0.28),
             font_size=Pt(8.5), bold=True, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

    row4_y = Inches(4.25)
    abh    = Inches(1.0)
    abw    = Inches(3.3)

    analyses = [
        (SCORE_CLR,
         "Confidence Scoring",
         "iPTM  ·  pLDDT  ·  Boltz-2 affinity\nAF3 ranking score  ·  PAE"),
        (VINA_CLR,
         "AutoDock Vina",
         "Local-only (minimized pose)\nRedocking (SMILES conformer)"),
        (KD_CLR,
         "Experimental Validation",
         "Spearman / Pearson correlation\nvs Kd  (n = 13 compounds, SPR)"),
    ]

    ana_centers = []
    a_gaps = [Inches(0.4), Inches(4.0), Inches(7.65)]
    for (clr, title, sub), gx in zip(analyses, a_gaps):
        box = add_rounded_rect(slide, gx, row4_y, abw, abh, fill=clr)
        add_label_in_shape(slide, box, f"{title}\n{sub}", font_size=Pt(8.5))
        ana_centers.append(gx + abw // 2)

    # Arrow summary → analyses
    sum_bot_y = row3_y + Inches(0.65)
    for ax in ana_centers:
        add_arrow(slide, sum_top_x, sum_bot_y, ax, row4_y, color=ARROW_CLR)

    # ── Row 5: Contact analysis side note ─────────────────────────────────
    add_text(slide,
             "★  Contact Analysis (gemmi distance search)  →  hot-spot residues contacted by multiple ligands",
             Inches(0.4), Inches(5.45), Inches(12), Inches(0.35),
             font_size=Pt(8.5), color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

    # ── Method legend dots ─────────────────────────────────────────────────
    legend_y = Inches(6.9)
    for i, (clr, lbl) in enumerate([(BOLTZ_CLR, "Boltz-2"),
                                     (AF3A_CLR,  "AF3 + AF3-MSA"),
                                     (AF3B_CLR,  "AF3 + ColabFold-MSA")]):
        lx = Inches(0.5 + i * 3.2)
        dot = add_rounded_rect(slide, lx, legend_y, Inches(0.18), Inches(0.18), fill=clr)
        add_text(slide, lbl, lx + Inches(0.25), legend_y - Inches(0.02),
                 Inches(2.5), Inches(0.25), font_size=Pt(8.5), color=LIGHT_GRAY,
                 align=PP_ALIGN.LEFT)

    return slide


# ── Slide 3: Confidence Rankings ──────────────────────────────────────────

def slide_confidence_rankings(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, DARK_BG)
    slide_header_footer(
        slide,
        "Confidence Rankings — All 3 Methods",
        "Boltz-2 iPTM  ·  AF3 iPTM (both conditions)  ·  Boltz-2 predicted affinity"
    )

    img_path = Path("ppt/fig_confidence_ranking.png")
    embed_figure(slide, img_path,
                 Inches(0.3), Inches(1.05),
                 Inches(12.73), Inches(5.75))

    add_caption(slide,
        "Left: dashed lines = sialic acid (control) iPTM per method. "
        "Right: Boltz-2 affinity value (negative = predicted tighter binder). "
        "Top Boltz-2 hits: cpd_4, cpd_1, cpd_9.  AF3 top hit: cpd_10, sa.")
    return slide


# ── Slide 4: AF3 Structural Divergence ────────────────────────────────────

def slide_af3_divergence(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, DARK_BG)
    slide_header_footer(
        slide,
        "AF3 Structural Divergence: Same Score, Different Structure",
        "AF3-MSA vs AF3-ColabFold  |  Cα RMSD per ligand"
    )

    img_path = Path("ppt/fig_af3_divergence.png")
    embed_figure(slide, img_path,
                 Inches(0.3), Inches(1.05),
                 Inches(8.5), Inches(4.6))

    # Bullet text box on the right
    bullets = [
        ("iPTM scores nearly identical (mean |Δ| = 0.034)",
         "similar confidence numbers hide structural disagreement"),
        ("Cα RMSD up to 77.9 Å — structurally completely different models",
         "7 of 14 ligands show RMSD > 30 Å"),
        ("Conclusion: AF3 is confidently inconsistent",
         "two runs yield same score, different structure"),
    ]

    bx = Inches(9.0)
    by = Inches(1.3)
    txb = slide.shapes.add_textbox(bx, by, Inches(4.1), Inches(5.0))
    tf = txb.text_frame
    tf.word_wrap = True

    first = True
    for header, body in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(10)

        run = p.add_run()
        run.text = "• " + header
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = BOLTZ_CLR

        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = "  " + body
        run2.font.size = Pt(10)
        run2.font.bold = False
        run2.font.color.rgb = LIGHT_GRAY

    return slide


# ── Slide 5: Contact Heatmap ──────────────────────────────────────────────

def slide_contact_heatmap(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, DARK_BG)
    slide_header_footer(
        slide,
        "Binding Site Hot-Spots — Contact Analysis",
        "Distance-based contacts (gemmi)  ·  Residues contacted across all 3 methods"
    )

    img_path = Path("ppt/contact_heatmap.png")
    embed_figure(slide, img_path,
                 Inches(0.3), Inches(1.05),
                 Inches(8.5), Inches(5.3))

    # Parse first 6 bullets from contact_bullets.txt
    bullets_raw = Path("ppt/contact_bullets.txt").read_text().splitlines()
    bullet_lines = []
    for line in bullets_raw[2:]:   # skip header lines
        stripped = line.strip()
        if stripped.startswith("●"):
            bullet_lines.append(stripped)
        if len(bullet_lines) == 6:
            break

    bx = Inches(9.0)
    by = Inches(1.3)
    txb = slide.shapes.add_textbox(bx, by, Inches(4.1), Inches(5.6))
    tf = txb.text_frame
    tf.word_wrap = True

    first = True
    for bullet in bullet_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(6)

        run = p.add_run()
        run.text = bullet[:120]   # truncate very long lines
        run.font.size = Pt(9.5)
        run.font.color.rgb = LIGHT_GRAY

    return slide


# ── Slide 6: Scores vs Experimental Kd ───────────────────────────────────

def slide_scatter_kd(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, DARK_BG)
    slide_header_footer(
        slide,
        "Computational Scores vs Experimental Kd  (n = 13)",
        "pKd = 6 − log₁₀(Kd µM)  ·  Higher pKd = tighter binding"
    )

    img_path = Path("ppt/fig_scatter_vs_kd.png")
    embed_figure(slide, img_path,
                 Inches(0.2), Inches(1.0),
                 Inches(12.93), Inches(6.15))

    return slide


# ── Slide 8: Summary & Top Hits ──────────────────────────────────────────

def slide_summary(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, SW, Inches(0.08), fill=BOLTZ_CLR)
    add_rect(slide, 0, SH - Inches(0.08), SW, Inches(0.08), fill=AF3A_CLR)

    add_text(slide, "Summary & Top Hits",
             Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.60),
             font_size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Two column layout
    sections = [
        # (header_color, header_text, body_lines)
        (BOLTZ_CLR,
         "Top Computational Hits — convergent across methods",
         ["cpd_3, cpd_6  (consistent mid-upper tier in Boltz-2 and both AF3 conditions)"]),

        (AF3A_CLR,
         "Boltz-2-specific top hits",
         ["cpd_4  (#1 by iPTM, 0.963)",
          "cpd_1  (#2 by iPTM, 0.946)"]),

        (AF3B_CLR,
         "AF3-specific top hits",
         ["cpd_10  (#1–2 in both AF3-MSA and AF3-ColabFold conditions)",
          "sa (sialic acid)  ranked highly by AF3, but contacts distinct site"]),

        (SCORE_CLR,
         "Consistent binding site residues (all 3 methods)",
         ["LYS124, TYR130, GLY131, TYR132"]),

        (VINA_CLR,
         "Recommended for wet-lab follow-up",
         ["cpd_3, cpd_9, cpd_12  —  tight experimental Kd (4.7–12.8 µM) + reasonable computational scores"]),

        (KD_CLR,
         "Key findings",
         ["Sialic acid contacts a distinct sub-pocket — synthetic analogs may target a different site",
          "AF3 pose quality assessed by Vina local-only (minimized energy of predicted pose)",
          "All 13 compounds show divergent binding modes across methods (Jaccard < 0.2)"]),
    ]

    # Column 1: sections 0-2, Column 2: sections 3-5
    col1_x = Inches(0.4)
    col2_x = Inches(6.9)
    col_w  = Inches(6.2)

    for col_idx in range(2):
        x = col1_x if col_idx == 0 else col2_x
        y_cursor = Inches(0.9)

        for sec_idx in range(3):
            global_idx = col_idx * 3 + sec_idx
            if global_idx >= len(sections):
                break
            hdr_clr, hdr_text, body_lines = sections[global_idx]

            # Section header
            txb = slide.shapes.add_textbox(x, y_cursor, col_w, Inches(0.38))
            tf = txb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = hdr_text
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = hdr_clr
            y_cursor += Inches(0.38)

            # Body lines
            for line in body_lines:
                txb2 = slide.shapes.add_textbox(x + Inches(0.2), y_cursor, col_w - Inches(0.2), Inches(0.30))
                tf2 = txb2.text_frame
                tf2.word_wrap = True
                p2 = tf2.paragraphs[0]
                run2 = p2.add_run()
                run2.text = "• " + line
                run2.font.size = Pt(10.5)
                run2.font.color.rgb = WHITE
                y_cursor += Inches(0.28)

            y_cursor += Inches(0.12)   # gap between sections

    return slide


# ── Slide 9: Methods ─────────────────────────────────────────────────────

def slide_methods(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, SW, Inches(0.08), fill=BOLTZ_CLR)
    add_rect(slide, 0, SH - Inches(0.08), SW, Inches(0.08), fill=AF3A_CLR)

    add_text(slide, "Methods",
             Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.60),
             font_size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    methods_data = [
        (BOLTZ_CLR, "Cofolding — Boltz-2",
         ["Version: Boltz-2 (open-source, RCSB pre-release)",
          "Compute: AWS EC2 g5.xlarge (NVIDIA A10G GPU, 24 GB VRAM)",
          "Runtime: ~8 min/ligand  ·  Docker container",
          "Input: FASTA protein + SMILES ligand + ColabFold .a3m MSA"]),

        (AF3A_CLR, "Cofolding — AlphaFold3",
         ["Version: AF3 web server (alphafoldserver.com), May 2025",
          "Condition 1 (AF3-MSA): AF3 built-in MSA pipeline",
          "Condition 2 (AF3-ColabFold): pre-computed MSA from ColabFold (same as Boltz-2)",
          "42 jobs total (14 ligands × 3 methods)"]),

        (SCORE_CLR, "MSA / Sequence",
         ["ColabFold MSA server (mmseqs2): pre-computed .a3m for Siglec-6",
          "FASTA: Siglec-6 Vtype + C2type1 + C2type2  (333 residues)"]),

        (VINA_CLR, "Docking — AutoDock Vina",
         ["Local-only: predicted pose locally minimized (energy relaxation, no global search)",
          "Redock: SMILES conformer (RDKit ETKDGv3) docked into predicted pocket",
          "Box size: 30 × 30 × 30 Å centered on predicted ligand centroid",
          "Receptor prep: remove ligand from CIF, convert with Open Babel"]),

        (KD_CLR, "Experimental Data",
         ["Kd values: SPR binding assay (n=13 compounds)",
          "Source: Hackathon dataset (Siglec-6 In Silico Library)",
          "No Kd for sialic acid (control excluded from correlation analysis)"]),

        (ARROW_CLR, "Contact Analysis",
         ["Tool: gemmi (Python)  ·  distance cutoff 5 Å",
          "Residue-ligand contacts extracted from cofolded CIF/mmCIF structures",
          "Jaccard similarity computed for top-5 contacted residues per compound"]),
    ]

    col1_x = Inches(0.4)
    col2_x = Inches(6.9)
    col_w  = Inches(6.2)

    for col_idx in range(2):
        x = col1_x if col_idx == 0 else col2_x
        y_cursor = Inches(0.9)

        for sec_idx in range(3):
            global_idx = col_idx * 3 + sec_idx
            if global_idx >= len(methods_data):
                break
            hdr_clr, hdr_text, body_lines = methods_data[global_idx]

            txb = slide.shapes.add_textbox(x, y_cursor, col_w, Inches(0.36))
            tf = txb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = hdr_text
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = hdr_clr
            y_cursor += Inches(0.36)

            for line in body_lines:
                txb2 = slide.shapes.add_textbox(x + Inches(0.2), y_cursor,
                                                col_w - Inches(0.2), Inches(0.28))
                tf2 = txb2.text_frame
                tf2.word_wrap = True
                p2 = tf2.paragraphs[0]
                run2 = p2.add_run()
                run2.text = "• " + line
                run2.font.size = Pt(9.5)
                run2.font.color.rgb = WHITE
                y_cursor += Inches(0.27)

            y_cursor += Inches(0.10)

    return slide


# ── Slide 4: Binding Site Overview ───────────────────────────────────────

def slide_binding_overview(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, DARK_BG)
    slide_header_footer(
        slide,
        "Where Do the Ligands Bind? — All 14 Compounds Overlaid",
        "Receptors aligned by Cα  ·  Each sphere = ligand center of mass  ·  Color = iPTM confidence"
    )

    img_path = Path("ppt/fig_binding_overview.png")
    embed_figure(slide, img_path,
                 Inches(0.2), Inches(1.0),
                 Inches(12.93), Inches(6.1))

    add_caption(slide,
        "Faint trace = pocket residues within 20 Å of any ligand.  "
        "Tight cluster → consistent binding site prediction.  "
        "Bold labels = top 3 compounds by iPTM confidence per method.")
    return slide


# ── Build full deck ────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    slide_title(prs)              # 1
    slide_pipeline(prs)           # 2
    slide_confidence_rankings(prs) # 3
    slide_binding_overview(prs)   # 4
    slide_af3_divergence(prs)     # 5
    slide_contact_heatmap(prs)    # 6
    slide_scatter_kd(prs)         # 7
    slide_summary(prs)            # 8
    slide_methods(prs)            # 9

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
